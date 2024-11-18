from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import pypandoc
import ocrmypdf
import json
import os
import warnings
warnings.filterwarnings("ignore", category=UserWarning)
# pypandoc.download_pandoc()


class PDF_Process():
    def __init__(self, pdf_path: str) -> None:
        """
        This function takes the path to a PDF file as input and determines whether the PDF needs OCR.
        If the PDF contains text, OCR is not needed; otherwise, OCR is required.
        input:
            pdf_path: str
        output:
            None
        """
        self.pdf_path = pdf_path
        self.tmp_output = f'{self.pdf_path[:-4]}_output_tmp.pdf'
        self.req_ocr = False  # Default to False

        reader = PdfReader(self.pdf_path)
        # Check if the PDF needs OCR
        for page in reader.pages:
            text = page.extract_text()
            if not text or not text.strip():
                # If any page does not contain text, OCR is needed
                self.req_ocr = True
                break  # No need to check further pages

        self._pdf_to_txt()
        pass

    def _pdf_to_txt(self) -> None:
        if self.req_ocr:
            ocrmypdf.ocr(self.pdf_path, self.tmp_output, force_ocr=True) #disable force_ocr
            self.process_text = self._read_pdf(self.tmp_output)
            if os.path.exists(self.tmp_output):
                os.remove(self.tmp_output)
        else:
            self.process_text = self._read_pdf(self.pdf_path)
        pass
    
    def _read_pdf(self, tmp_pdf_path: str) -> str:
        try:
            tmp_text = ''
            reader = PdfReader(tmp_pdf_path)
            for page in reader.pages:
                tmp_text += page.extract_text()
        
            extracted_text = tmp_text
            
            # Cleanup the extracted text
            extracted_text = extracted_text.strip()
            extracted_text = os.linesep.join([s for s in extracted_text.splitlines() if s.strip()])

            return extracted_text
        except:
            return ""

# example usage
# extracted_doc = PDF_Process(pdf_path='before_ocr.pdf').process_text #get_processed_text())

class DOCX_Process():
    def __init__(self, docx_file):
        '''
        This function initializes the DocxToTextConverter object with the provided file path.
        init method
            input: file path
            output: None
        '''
        self.file_path = docx_file
        self.extracted_text = self.convert_to_text().strip()

    def convert_to_text(self):
        '''
        This function reads the .docx file and converts its content to text using the docx library.
        convert_to_text method
            input: None
            output: text content of the .docx file
        '''
        try:
            doc = Document(self.file_path)
            text = '\n'.join([para.text for para in doc.paragraphs])
            if '[TABLE]' in text:
                text += self._extract_tables()
            return text
        except Exception as e:
            print(f"An error occurred: {e}")
            return ''

    def _extract_tables(self):
        '''
        This function extracts tables from the .docx file and converts them to text.
        _extract_tables method
            input: None
            output: text content of the tables in the .docx file
        '''
        try:
            tables_text = ''
            doc = Document(self.file_path)
            for table in doc.tables:
                print(f"Table found: {table}")
                for row in table.rows:
                    for cell in row.cells:
                        tables_text += cell.text + '\t'
                    tables_text += '\n'
                tables_text += '\n'
            return tables_text
        except Exception as e:
            print(f"An error occurred while extracting tables: {e}")
            return ''

class PPTX_Process():
    def __init__(self, file_path):
        self.file_path = file_path
        self.slides_content = self.read_pptx()
        self.extracted_text = self.combine_slides_text()

    def read_pptx(self):
        """
        This function reads a PowerPoint file and extracts the text from each slide.
        
        Args:
        file_path (str): The path to the PowerPoint file.
        
        Returns:
        list: A list of strings, each representing the text content of a slide.
        """
        prs = Presentation(self.file_path)
        slides_text = []

        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slides_text.append("\n".join(slide_text))
        
        return slides_text

    def combine_slides_text(self):
        """
        This function combines the text from all slides into a single string.
        
        Returns:
        str: The combined text content of all slides.
        """
        combined_text = ''
        for i, slide in enumerate(self.slides_content):
            combined_text += f'Slide {i+1}\n{slide}\n\n'
        return combined_text

def load_and_flatten_urls(json_path):
    """
    Load URLs from a JSON file and flatten the structure.

    Args:
    json_path (str): The path to the JSON file containing the URLs.

    Returns:
    dict: A dictionary with project IDs as keys and file URLs as values.
    """
    with open(json_path, 'r') as f:
        urls_data = json.load(f)

    flattened_urls_data = {}

    for project_id, data in urls_data.items():
        flattened_urls_data[project_id] = {}
        if 'Files' in data:
            for file_name, file_url in data['Files'].items():
                flattened_urls_data[project_id][file_name] = file_url
        if 'Folders' in data:
            for folder, files in data['Folders'].items():
                for file_name, file_url in files['Files'].items():
                    flattened_urls_data[project_id][file_name] = file_url

    return flattened_urls_data